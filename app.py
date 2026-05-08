from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
import easyocr
import numpy as np
from PIL import Image
import io
import fitz  # pymupdf
from pptx import Presentation
from groq import Groq

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

reader = easyocr.Reader(['ko', 'en'])
client = Groq(api_key="GROQ_API_KEY")

def extract_from_image(contents):
    image = Image.open(io.BytesIO(contents))
    img_array = np.array(image)
    results = reader.readtext(img_array)
    return "\n".join([text for (_, text, conf) in results if conf > 0.3])

def extract_from_pdf(contents):
    doc = fitz.open(stream=contents, filetype="pdf")
    extracted = ""
    for page in doc:
        extracted += page.get_text()
    return extracted

def extract_from_pptx(contents):
    prs = Presentation(io.BytesIO(contents))
    extracted = ""
    for i, slide in enumerate(prs.slides):
        extracted += f"\n[슬라이드 {i+1}]\n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                extracted += shape.text_frame.text + "\n"
    return extracted

@app.post("/predict")
async def predict(file: UploadFile = File(...)):
    contents = await file.read()
    ext = file.filename.split(".")[-1].lower()

    # 확장자별 텍스트 추출
    if ext in ["jpg", "jpeg", "png", "webp"]:
        extracted = extract_from_image(contents)
    elif ext == "pdf":
        extracted = extract_from_pdf(contents)
    elif ext in ["ppt", "pptx"]:
        extracted = extract_from_pptx(contents)
    else:
        return {"result": "지원하지 않는 파일 형식이에요. (이미지/PDF/PPT만 가능)", "raw": ""}

    if not extracted.strip():
        return {"result": "텍스트를 추출할 수 없었어요.", "raw": ""}

    # LLM으로 정리
    response = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[
            {
                "role": "user",
                "content": f"당신은 문서 정리 전문가입니다. 다음 데이터는 문서에서 추출한 내용이며, 해당 문서의 종류를 추측하여 알아보기 쉽게 정리하고, 문서의 내용을 요약해주세요.:\n\n{extracted}"
            }
        ]
    )

    summary = response.choices[0].message.content
    return {"result": summary, "raw": extracted}